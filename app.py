import os
from flask import Flask, render_template, request, send_from_directory
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image

app = Flask(__name__)

# Konfiguration: Ordner für das Hochladen von Dateien und das Speichern von Logos
UPLOAD_FOLDER = 'uploads'
LOGO_FOLDER = 'static/logos'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['LOGO_FOLDER'] = LOGO_FOLDER

# Sicherstellen, dass die Ordner existieren
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(LOGO_FOLDER, exist_ok=True)


@app.route('/')
def index():
    return render_template('index.html')  # Lade das HTML-Upload-Formular


@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return 'Keine Datei ausgewählt', 400
    file = request.files['file']

    if file.filename == '':
        return 'Keine Datei ausgewählt', 400

    # Speichere die Datei im Upload-Ordner
    filename = file.filename
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)

    # Bearbeite die PPTX-Datei
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"modified_{filename}")
    logo_path = os.path.join(app.config['LOGO_FOLDER'], 'mdrfragtlogo.jpg')  # Dein Logo-Pfad

    modify_presentation(file_path, output_path, logo_path)  # Funktion zum Bearbeiten der PPTX

    return send_from_directory(app.config['UPLOAD_FOLDER'], f"modified_{filename}", as_attachment=True)


def modify_presentation(input_pptx, output_pptx, logo_path):
    prs = Presentation(input_pptx)

    # Farbwerte und Logo
    background_color = RGBColor(174, 177, 192)
    bottom_color = RGBColor(255, 105, 86)
    font_color = RGBColor(0, 0, 0)

    chart_colors = [
        RGBColor(255, 105, 86),  # Rot
        RGBColor(232, 230, 215),  # Beige
        RGBColor(255, 255, 255),  # Weiß
        RGBColor(0, 0, 0)  # Schwarz
    ]

    # Designänderungen auf alle Folien anwenden
    apply_design_modifications(prs, background_color, bottom_color, logo_path)

    # Durch alle Folien iterieren und Texte und Diagramme anpassen
    for slide in prs.slides:
        style_textbox(slide, RGBColor(255, 105, 86))
        add_bullet_textbox(prs, slide, ["Erster Punkt", "Zweiter Punkt"])

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                chart_type = chart.chart_type

                # Anpassung von Balkendiagrammen
                if chart_type in [XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.BAR_CLUSTERED]:
                    modify_bar_chart(chart, shape, prs, font_color, chart_colors)

                # Anpassung von Tortendiagrammen
                elif chart_type == XL_CHART_TYPE.PIE:
                    modify_pie_chart(chart, shape, prs, font_color, chart_colors)

    # Präsentation speichern
    prs.save(output_pptx)


def apply_design_modifications(prs, background_color, bottom_color, logo_path):
    for slide in prs.slides:
        set_background_color(slide, background_color)
        color_bottom_slide(prs, slide, bottom_color)
        add_logo(prs, slide, logo_path)


def set_background_color(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def color_bottom_slide(prs, slide, color):
    slide_height = prs.slide_height
    height = slide_height // 20  # Bottom 1/20th of the slide
    shape = slide.shapes.add_shape(
        1, Inches(0), slide_height - height, prs.slide_width, height  # 1 corresponds to MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color  # No border color


def add_logo(prs, slide, logo_path):
    left = prs.slide_width - Inches(1.9)
    top = Inches(0.0)

    # Berechnung der Höhe und Breite des Logos
    with Image.open(logo_path) as img:
        aspect_ratio = img.width / img.height
    height = Inches(1)
    width = height * aspect_ratio

    slide.shapes.add_picture(logo_path, left, top, width=width, height=height)


def add_bullet_textbox(prs, slide, text_list, font_name="Frutiger 44 Light", font_size=Pt(16),
                       font_color=RGBColor(0, 0, 0)):
    left = Inches(7.5)
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Text oben ausrichten

    for i, text in enumerate(text_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {text}"
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT


def style_textbox(slide, color):
    text_boxes = [s for s in slide.shapes if s.has_text_frame]
    if len(text_boxes) > 0:
        text_boxes[0].text_frame.clear()
        paragraph = text_boxes[0].text_frame.add_paragraph()
        run_title = paragraph.add_run()
        run_title.text = 'TITEL EINFÜGEN'
        run_title.font.size = Pt(20)
        run_title.font.bold = True
        run_title.font.name = "Frutiger 44 Light"

    if len(text_boxes) > 1:
        tb = text_boxes[1]
        tb.line.color.rgb = color
        tb.line.width = Pt(1)
        tb.height = tb.height // 2
        existing_text = tb.text_frame.text.strip()
        tb.text_frame.clear()
        p = tb.text_frame.add_paragraph()
        run_bold = p.add_run()
        run_bold.text = "Frage: "
        run_bold.font.bold = True
        run_bold.font.size = Pt(10)
        run_bold.font.name = "Frutiger 44 Light"

        run_normal = p.add_run()
        run_normal.text = existing_text
        run_normal.font.bold = False
        run_normal.font.size = Pt(10)
        run_normal.font.name = "Frutiger 44 Light"


def modify_bar_chart(chart, shape, prs, font_color, colors):
    for i, series in enumerate(chart.plots[0].series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]
        for point in series.points:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]

    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 1
    shape.top = int(prs.slide_height / 5)
    shape.width = int(prs.slide_width / 2)


def modify_pie_chart(chart, shape, prs, font_color, colors):
    for i, point in enumerate(chart.plots[0].series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]

    shape.left = 0
    shape.top = int(prs.slide_height / 5)
    shape.width = int(prs.slide_width / 2)


if __name__ == '__main__':
    app.run(debug=True)
