import os
print("Aktuelles Arbeitsverzeichnis:", os.getcwd())
print("Datei gefunden?", os.path.isfile("input.pptx"))
print("Dateien im Ordner:", os.listdir())
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_MARK, XL_TICK_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def add_bullet_textbox(prs, slide, text_list, font_name="Frutiger 44 Light", font_size=Pt(16),
                       font_color=RGBColor(0, 0, 0)):
    """
    Fügt eine TextBox mit Spiegelstrichen auf der rechten Hälfte der Folie hinzu.

    :param slide: Die Folie, auf die die TextBox eingefügt wird.
    :param text_list: Eine Liste von Strings, die als Spiegelstriche hinzugefügt werden.
    :param font_name: Die gewünschte Schriftart (Standard: Frutiger 44 Light).
    :param font_size: Die Schriftgröße (Standard: 16pt).
    :param font_color: Die Schriftfarbe (Standard: Schwarz).
    """
    # Position der TextBox (rechte Hälfte der Folie)
    left = Inches(7.5)  # Start auf der rechten Seite der Folie

    top = Inches(1.5)  # Obere Position
    width = Inches(5)  # Breite der TextBox
    height = Inches(5)  # Höhe der TextBox

    # TextBox hinzufügen
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Text oben ausrichten

    # Ersten Absatz hinzufügen (erster Spiegelstrich)
    for i, text in enumerate(text_list):
        if i == 0:
            p = text_frame.paragraphs[0]  # Erster Absatz
        else:
            p = text_frame.add_paragraph()  # Neue Zeile für weitere Spiegelstriche
        p.text = f"• {text}"  # Spiegelstrich hinzufügen
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = True  # **Text fett setzen**

        p.alignment = PP_ALIGN.LEFT  # Links ausrichten


# Set the background color of a slide
def set_background_color(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# Add a colored rectangle at the bottom of a slide
def color_bottom_slide(prs, slide, color):
    slide_height = prs.slide_height
    height = slide_height // 20  # Bottom 1/20th of the slide
    shape = slide.shapes.add_shape(
        1, Inches(0), slide_height - height, prs.slide_width, height  # 1 corresponds to MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color  # No border color


def add_logo(prs, slide, logo_path):
    # Annahme: Die Überschrift ist das erste Textfeld auf der Folie
    headline_height = None
    for shape in slide.shapes:
        if shape.has_text_frame:  # Überprüfen, ob das Shape ein Textfeld ist
            headline_height = shape.height  # Höhe des Textfelds speichern
            break  # Schleife beenden, sobald das erste Textfeld gefunden wurde

    # Berechnung der Position für das Logo
    left = prs.slide_width - Inches(
        1.9)  # Platzierung des Logos 1,9 Zoll von der rechten Kante(kante schließt mit textbox ab)
    top = Inches(0.0)  # Verschiebe das Logo 0 Zoll von der oberen Kante, damit sie ganz oben ist

    # Berechnung der Höhe und Breite des Logos
    if headline_height is not None:
        height = headline_height * 1.5  # Vergrößere das Logo um 50% der Überschriftenhöhe
    else:
        height = Inches(1)  # Standardhöhe auf 1 Zoll setzen, wenn keine Überschrift gefunden wird

    # Proportionale Breite berechnen, um Verzerrungen zu vermeiden
    from PIL import Image
    with Image.open(logo_path) as img:
        aspect_ratio = img.width / img.height  # Seitenverhältnis des Logos berechnen
    width = height * aspect_ratio  # Breite basierend auf der Höhe berechnen

    # Logo auf der Folie platzieren
    slide.shapes.add_picture(logo_path, left, top, width=width, height=height)


# Apply design modifications to all slides in the presentation
def apply_design_modifications(prs, background_color, bottom_color, logo_path):
    for slide in prs.slides:
        set_background_color(slide, background_color)
        color_bottom_slide(prs, slide, bottom_color)
        add_logo(prs, slide, logo_path)


# Modify colors of a bar chart
def modify_bar_chart_colors(chart, colors):
    for i, series in enumerate(chart.plots[0].series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]
        for point in series.points:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]


from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.util import Pt
from pptx.dml.color import RGBColor


def style_textbox(slide, color):
    # Alle Textfelder (Textboxen) auf der Folie erfassen
    text_boxes = [s for s in slide.shapes if s.has_text_frame]

    # Erstes Textfeld bearbeiten (Titel)
    if len(text_boxes) > 0:
        text_boxes[0].text_frame.clear()  # Vorhandenen Inhalt löschen
        paragraph = text_boxes[0].text_frame.add_paragraph()  # Neuen Absatz hinzufügen
        run_title = paragraph.add_run()  # Neuen Textlauf hinzufügen
        run_title.text = 'TITEL EINFÜGEN'  # Standardtext setzen
        run_title.font.size = Pt(20)  # Schriftgröße auf 20 Punkte setzen
        run_title.font.bold = True  # Fett formatieren
        run_title.font.name = "Frutiger 44 Light"  # Schriftart setzen

    # Zweites Textfeld bearbeiten (Frage)
    if len(text_boxes) > 1:
        tb = text_boxes[1]  # Das zweite Textfeld abrufen
        tb.line.color.rgb = color  # Rahmenfarbe setzen
        tb.line.width = Pt(1)  # Rahmenbreite setzen
        tb.height = tb.height // 2  # Höhe halbieren

        # Vorhandenen Text beibehalten
        existing_text = tb.text_frame.text.strip()  # Ursprünglichen Text sichern

        # Textfeld leeren, um neuen Absatz mit Formatierungen hinzuzufügen
        tb.text_frame.clear()
        p = tb.text_frame.paragraphs[0] if tb.text_frame.paragraphs else tb.text_frame.add_paragraph()

        # p = tb.text_frame.add_paragraph()  # Neuen Absatz hinzufügen

        # "Frage: " hinzufügen und fett formatieren
        run_bold = p.add_run()
        run_bold.text = "Frage: "  # Text hinzufügen
        run_bold.font.bold = True  # Fettgedruckt setzen
        run_bold.font.size = Pt(10)  # Schriftgröße 10 Pt setzen
        run_bold.font.name = "Frutiger 44 Light"  # Schriftart setzen

        # Restlichen Text hinzufügen (nicht fett)
        run_normal = p.add_run()
        run_normal.text = existing_text  # Originaltext hinzufügen
        run_normal.font.bold = False  # Kein Fett
        run_normal.font.size = Pt(10)  # Schriftgröße 10 Pt setzen
        run_normal.font.name = "Frutiger 44 Light"  # Schriftart setzen

        # Kontrolle des Endergebnisses
        print(tb.text_frame.text)


# Modify colors of a pie chart
def modify_pie_chart_colors(chart, colors):
    for i, point in enumerate(chart.plots[0].series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]


# Funktion zur Anpassung von Balkendiagrammen
# Funktion zur Anpassung von Balkendiagrammen
def modify_bar_chart(chart, shape, prs, font_color, colors):
    # Farben für die Datenreihen im Diagramm festlegen
    modify_bar_chart_colors(chart, colors)

    # Die Werteachse (Y-Achse) wird von 0 bis 1 skaliert
    chart.value_axis.minimum_scale = 0  # Mindestwert der Achse auf 0 setzen
    chart.value_axis.maximum_scale = 1  # Höchstwert der Achse auf 1 setzen

    # Positionierung des Tortendiagramms auf der Folie
    # shape.left = int(prs.slide_width/ 20)
    shape.top = int(prs.slide_height / 5)  # Diagramm etwas nach unten verschieben
    shape.width = int(prs.slide_width / 2)  # Breite des Diagramms mit Rand


#  shape.height = int(prs.slide_height )  # Höhe des Diagramms auf die halbe Folie setzen

# Apply modifications to pie charts
def modify_pie_chart(chart, shape, prs, font_color, colors):
    # Farben für das Tortendiagramm anpassen
    modify_pie_chart_colors(chart, colors)

    # Positionierung des Tortendiagramms auf der Folie
    shape.left = 0  # Diagramm in der linken Hälfte positionieren
    shape.top = int(prs.slide_height / 5)  # Diagramm etwas nach unten verschieben
    shape.width = int(prs.slide_width / 2)  # Breite des Diagramms mit Rand
    # shape.height = int(prs.slide_height / 2)  # Höhe des Diagramms auf die halbe Folie setzen


# Main function to modify the presentation
def modify_presentation(input_pptx, output_pptx, logo_path):
    # Präsentation öffnen
    prs = Presentation(input_pptx)

    # Definition der Farbwerte für das Design
    background_color = RGBColor(174, 177, 192)  # Hintergrundfarbe
    bottom_color = RGBColor(255, 105, 86)  # Farbe für die untere Leiste
    font_color = RGBColor(0, 0, 0)  # Schriftfarbe (schwarz)

    # Farben für Diagramme festlegen (Reihenfolge geändert)
    chart_colors = [
        RGBColor(255, 105, 86),  # Rot
        RGBColor(232, 230, 215),  # Beige
        RGBColor(255, 255, 255),  # Weiß
        RGBColor(0, 0, 0)  # Schwarz
    ]

    # Designänderungen auf die Präsentation anwenden (Hintergrund, Logo, Farben)
    apply_design_modifications(prs, background_color, bottom_color, logo_path)

    # Seitennummerierung hinzufügen
    # add_slide_numbers(prs)

    # Durch alle Folien der Präsentation iterieren
    for slide in prs.slides:
        prs.core_properties.slides = True  # Aktiviert die Foliennummerierung in den Eigenschaften

        # Falls das Shape Text enthält, dann dessen Formatierung anpassen
        style_textbox(slide, RGBColor(255, 105, 86))
        text_items = ["Erster Punkt", "Zweiter Punkt"]
        add_bullet_textbox(prs, slide, text_items)

        for shape in slide.shapes:
            #
            # Falls das Shape ein Diagramm enthält, dieses formatieren
            if shape.has_chart:
                chart = shape.chart
                chart_type = chart.chart_type

                # Überprüfung, ob das Diagramm ein Balken- oder Säulendiagramm ist
                if chart_type in [
                    XL_CHART_TYPE.COLUMN_CLUSTERED,  # Gruppierte Säulen
                    XL_CHART_TYPE.COLUMN_STACKED,  # Gestapelte Säulen
                    XL_CHART_TYPE.COLUMN_STACKED_100,  # Gestapelte 100%-Säulen
                    XL_CHART_TYPE.BAR_CLUSTERED,  # Gruppierte Balken
                    XL_CHART_TYPE.BAR_STACKED,  # Gestapelte Balken
                    XL_CHART_TYPE.BAR_STACKED_100  # Gestapelte 100%-Balken
                ]:
                    # Anpassung der Balkendiagramme
                    modify_bar_chart(chart, shape, prs, font_color, chart_colors)

                # Überprüfung, ob das Diagramm ein Tortendiagramm ist
                elif chart_type == XL_CHART_TYPE.PIE:
                    # Anpassung der Tortendiagramme
                    modify_pie_chart(chart, shape, prs, font_color, chart_colors)

        # Falls das Shape Text enthält, dann dessen Formatierung anpassen
    #        style_textbox(slide, RGBColor(255, 105, 86))
    # Präsentation mit den vorgenommenen Änderungen speichern
    prs.save(output_pptx)

if __name__ == "__main__":
    import os

    # Basisverzeichnis: eine Ebene über /utils/
    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    input_pptx = os.path.join(BASE_DIR, "input.pptx")
    output_pptx = os.path.join(BASE_DIR, "output.pptx")
    logo_path = os.path.join(BASE_DIR, "logo.jpg")

    modify_presentation(input_pptx, output_pptx, logo_path)